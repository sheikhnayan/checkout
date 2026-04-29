from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

IMG = 'C:/tmp/pptx_images'

PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)
SECONDARY = RGBColor(0x2E, 0x86, 0xAB)
ACCENT = RGBColor(0xE8, 0x6F, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x2D, 0x3A, 0x4A)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox


def add_para(slide, left, top, width, height, text, size=16, color=DARK_TEXT, spacing=Pt(10)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.space_after = spacing
    p.line_spacing = 1.15
    return txBox


def add_img(slide, name, left, top, width, height):
    path = os.path.join(IMG, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)


# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_img(slide, 'teacher1.jpg', Inches(7.5), Inches(0), Inches(5.833), Inches(7.5))
add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.15), Inches(3.0), ACCENT)

add_text(slide, Inches(1.2), Inches(2.0), Inches(5.8), Inches(2.2),
         'Stuck in 18th Century',
         40, True, WHITE, font='Calibri Light')

add_text(slide, Inches(1.2), Inches(4.8), Inches(5), Inches(0.5), 'Koli Akter', 20, True, WHITE)
add_text(slide, Inches(1.2), Inches(5.3), Inches(5), Inches(0.4), 'ID: 056', 16, False, RGBColor(0xB0, 0xC4, 0xD8))
add_text(slide, Inches(1.2), Inches(5.7), Inches(5), Inches(0.4), 'Course Code: ENG 3101', 16, False, RGBColor(0xB0, 0xC4, 0xD8))
add_text(slide, Inches(1.2), Inches(6.1), Inches(5), Inches(0.4), 'Introduction To ELT', 16, False, RGBColor(0xB0, 0xC4, 0xD8))

# ===== SLIDE 2: Introduction =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Introduction', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'The 18th century is often called the Age of Enlightenment, a period when people began to rely more on science, logic, and reason rather than tradition or blind belief. It was a time of questioning—people started asking why things were the way they were and searched for answers through knowledge and observation.',
16, DARK_TEXT)

add_img(slide, 'teaching2.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 3: Today's World =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), PRIMARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), "Today's World", 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Today, we live in a completely different world, shaped by digital technology, the internet, and advanced science. However, even with all this progress, a question still remains: are we truly different from people of the 18th century, or are we simply a modern version of them?',
16, DARK_TEXT)

add_img(slide, 'books.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 4: Surface Changes =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Surface Changes', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'At first glance, it seems like we have changed a lot. Our world is faster, more connected, and more technologically advanced. We use smartphones, social media, artificial intelligence, and digital platforms in our daily lives. Information is available instantly, and communication happens across the globe within seconds. In this sense, we have clearly moved far beyond the 18th century.',
16, DARK_TEXT)

add_img(slide, 'students.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 5: Deeper Similarities =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'But if we look deeper...', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'But if we look deeper, the similarities become clearer. Like people of the Enlightenment, we strongly believe in science and rational thinking. We trust doctors, scientists, and technology. During global crises, such as pandemics, people turn to science for solutions.',
16, DARK_TEXT)

add_img(slide, 'graduation.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 6: Education & Knowledge =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), PRIMARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Education & Knowledge', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Education is highly valued, just as it was during the Enlightenment. Students are encouraged to think critically, analyze information, and form logical arguments. This shows that the foundation built in the 18th century is still very much alive today.',
16, DARK_TEXT)

add_img(slide, 'classroom.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 7: Enlightenment & Social Competition =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'The Other Side of Enlightenment', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'However, the Enlightenment period was not only about knowledge and reason. It was also a time of social competition, pride, and personal image. People cared deeply about how they were seen by others. Social status, fashion, and reputation played a huge role in everyday life. In many ways, society was built on appearance and public perception.',
16, DARK_TEXT)

add_img(slide, 'idea.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 8: Modern Digital Presentation =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Digital Presentation Today', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'In the modern digital world, especially on platforms like Instagram, people present carefully designed versions of their lives. Photos are edited, moments are selected, and captions are crafted to create a certain image. People often show their happiest, most successful, or most attractive sides, while hiding their struggles and imperfections. This creates a gap between reality and representation.',
16, DARK_TEXT)

add_img(slide, 'writing.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 9: The Perfect Day Illusion =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), PRIMARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'The Illusion of Perfection', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'For example, a person might post a picture of a perfect day out, smiling and looking confident. But what we do not see is the stress, the bad moments, or the ordinary parts of their life. This selective sharing creates an illusion, and others who see it may start comparing their own lives to this unrealistic standard.',
16, DARK_TEXT)

add_img(slide, 'teaching2.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 10: Not New, Just Modern =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'An Old Habit', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'This behavior is not new—it is simply a modern version of an old habit. In the 18th century, people showed their status through clothing, manners, and social gatherings. Today, people do the same thing through posts, followers, and likes. The platform has changed, but the intention remains the same. Humans still want to be admired, respected, and noticed.',
16, DARK_TEXT)

add_img(slide, 'books.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 11: Validation =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'The Issue of Validation', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Another important issue is validation. In the past, validation came from society—people wanted approval from their community or social circle. Today, validation comes in the form of likes, comments, and shares. A simple post can become a source of happiness or disappointment depending on how others react to it. This shows how deeply social approval still affects us.',
16, DARK_TEXT)

add_img(slide, 'students.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 12: Digital Pressure =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), PRIMARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Amplified Pressure', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'At the same time, the digital world has made this pressure even stronger. In the 18th century, social comparison was limited to a small group of people. Now, through social media, we compare ourselves with hundreds or even thousands of others. This can create anxiety, low self-esteem, and a constant feeling of not being "good enough."',
16, DARK_TEXT)

add_img(slide, 'graduation.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 13: Misuse of Knowledge =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Misuse of Knowledge', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Even though we have access to more knowledge than ever before, we do not always use it wisely. The Enlightenment thinkers believed that reason would lead to progress and a better society. However, today we often see the opposite. Misinformation spreads quickly online, people follow trends without thinking, and sometimes emotions are valued more than facts.',
16, DARK_TEXT)

add_img(slide, 'classroom.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 14: Examples of Misinformation =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Following Without Thinking', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'For example, many people share information on social media without checking if it is true. Others follow popular opinions just to fit in, rather than forming their own views. This shows that although we have the tools for critical thinking, we do not always apply them.',
16, DARK_TEXT)

add_img(slide, 'idea.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 15: Technology & Identity =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), PRIMARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Technology & Identity', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Another interesting point is how technology affects our identity. In the digital world, people can create a completely different version of themselves. They can choose what to show and what to hide. Over time, this can make it difficult to separate real identity from online identity.',
16, DARK_TEXT)

add_img(slide, 'writing.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 16: Self-Worth & Online Presence =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Self-Worth & Online Presence', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'In some cases, people begin to measure their self-worth based on their online presence. This can be dangerous, especially for young people, who are still developing their sense of self. The need to appear perfect can lead to stress and dissatisfaction.',
16, DARK_TEXT)

add_img(slide, 'teaching2.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 17: The Positive Side =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Positive Opportunities', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Despite all these similarities, it is important to recognize that there are also differences. Today, we have more opportunities for learning, expression, and connection. Social media can be used for education, awareness, and creativity. People can share ideas, support causes, and connect with others across cultures.',
16, DARK_TEXT)

add_img(slide, 'books.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 18: How We Use Technology =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), PRIMARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'The Problem Is Usage', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'So, the problem is not technology itself, but how we use it. If we use technology wisely, we can continue the positive legacy of the Enlightenment—spreading knowledge, encouraging critical thinking, and improving society. But if we focus only on image, validation, and superficial success, then we are simply repeating the same patterns as the past.',
16, DARK_TEXT)

add_img(slide, 'students.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 19: Conclusion =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Conclusion', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'In conclusion, it is clear that we are not completely different from people of the 18th century. While our world has changed in terms of technology and science, human nature remains largely the same. Our desire for recognition, our concern for image, and our need for approval continue to shape our behavior. The only difference is that these qualities now exist in a digital form.',
16, DARK_TEXT)

add_img(slide, 'graduation.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 20: Final Thoughts =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), SECONDARY)
add_rect(slide, Inches(1.0), Inches(0.5), Inches(1.5), Inches(0.06), ACCENT)
add_text(slide, Inches(1.0), Inches(0.7), Inches(10), Inches(0.8), 'Final Thoughts', 32, True, PRIMARY, font='Calibri Light')

add_para(slide, Inches(1.0), Inches(1.6), Inches(6), Inches(5.3),
'Therefore, it can be said that we are still "stuck in the 18th century"—not because we have failed to progress, but because we carry the same human tendencies into every new era. The challenge for us is not just to advance technologically, but to grow intellectually and emotionally as well.',
16, DARK_TEXT)

add_img(slide, 'classroom.jpg', Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.8))

# ===== SLIDE 21: Thank You =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

add_text(slide, Inches(2), Inches(2.0), Inches(9.333), Inches(1.5),
         'Thank You!', 54, True, WHITE, PP_ALIGN.CENTER, 'Calibri Light')
add_rect(slide, Inches(5.5), Inches(3.7), Inches(2.333), Inches(0.08), ACCENT)
add_text(slide, Inches(2), Inches(4.2), Inches(9.333), Inches(1),
         'Any Questions?', 28, False, RGBColor(0xB0, 0xC4, 0xD8), PP_ALIGN.CENTER, 'Calibri Light')
add_text(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(0.5),
         'Koli Akter  |  ID: 056  |  ENG 3101 - Introduction To ELT', 16, False, RGBColor(0x8A, 0xA0, 0xB8), PP_ALIGN.CENTER)

# Save
output = 'c:/wamp64/www/checkout/Stuck_in_18th_Century_Presentation.pptx'
prs.save(output)
print(f'Saved: {output}')
print(f'Total slides: {len(prs.slides)}')
